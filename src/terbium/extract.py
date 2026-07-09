"""Export the actual product images, not just the text.

This pulls every embedded photo out of a PDF, and - for lookbooks - names each
file after the product it sits under. It also reports the qualities that matter
downstream: pixel size, colorspace, format, position, effective DPI, dominant
colour, and the associated product name and collection.

Bytes come from PyMuPDF (already a terbium dependency), which hands back a
ready-to-save JPEG/PNG. Repeated images (a logo stamped on every page) are
detected and skipped so you do not export the same watermark 80 times.
"""
from __future__ import annotations

import io
import os
import re
from collections import defaultdict
from typing import List, Optional

import fitz

from .documents.pdf import PdfAdapter
from .layout import labels as _labels
from .layout.images import classify
from .layout.lines import cluster_lines
from .model.elements import ImageRef, Page

_CS = {1: "Gray", 3: "RGB", 4: "CMYK"}


def _safe(name: str) -> str:
    return re.sub(r"[^A-Za-z0-9._-]+", "_", name).strip("_")[:80] or "image"


def _label_positions(page: Page) -> List[dict]:
    """Merged product labels on a page, each with a centre and baseline."""
    cands = []
    for ln in cluster_lines(page.words):
        for seg in _labels._segment_row(ln):
            if _labels._is_label(seg):
                cands.append(seg)
    return _labels._merge_wrapped(cands) if cands else []


def _associate(bbox, items: List[dict]) -> Optional[str]:
    """The label sitting just below an image, horizontally aligned to it."""
    if not bbox or not items:
        return None
    ix0, iy0, ix1, iy1 = bbox
    best, best_dy = None, 1e9
    for it in items:
        if ix0 - 25 <= it["cx"] <= ix1 + 25 and it["y"] > iy1 - 8:
            dy = it["y"] - iy1
            if dy < best_dy and dy < 140:
                best, best_dy = it, dy
    return best["name"] if best else None


def _dominant_color(data: bytes) -> Optional[str]:
    try:
        from PIL import Image

        with Image.open(io.BytesIO(data)) as im:
            im = im.convert("RGB").resize((16, 16))
            px = list(im.getdata())
        r = sum(p[0] for p in px) // len(px)
        g = sum(p[1] for p in px) // len(px)
        b = sum(p[2] for p in px) // len(px)
        return f"#{r:02x}{g:02x}{b:02x}"
    except Exception:
        return None


def _dpi(width_px: int, bbox) -> Optional[int]:
    if not bbox:
        return None
    width_pt = bbox[2] - bbox[0]
    if width_pt <= 0:
        return None
    return round(width_px / (width_pt / 72.0))


def export_images(
    path: str,
    out_dir: str,
    only_photos: bool = True,
    associate: bool = True,
    dedupe_repeats: bool = True,
    min_side: int = 180,
    max_aspect: float = 4.0,
) -> List[dict]:
    """Extract images from ``path`` into ``out_dir``; return a manifest.

    ``only_photos`` skips icons, swatches, thin banners, and slivers, keeping real
    product photos (shorter side >= ``min_side`` px, aspect ratio <= ``max_aspect``).
    ``associate`` names lookbook photos after the product beneath them.
    ``dedupe_repeats`` drops logos/watermarks that recur across many pages.
    """
    os.makedirs(out_dir, exist_ok=True)
    pages = PdfAdapter().parse(path)
    n = max(1, len(pages))

    xref_pages = defaultdict(set)
    for p in pages:
        for im in p.images:
            if im.xref:
                xref_pages[im.xref].add(p.index)
    repeated = {x for x, pgs in xref_pages.items() if len(pgs) >= max(3, 0.2 * n)}

    manifest: List[dict] = []
    used_names: dict = {}
    doc = fitz.open(path)
    try:
        for p in pages:
            items = _label_positions(p) if associate else []
            collection = _labels._pick_collection(items, p) if items else None
            for im in p.images:
                if only_photos:
                    if im.kind != "photo":
                        continue
                    lo, hi = min(im.width, im.height), max(im.width, im.height)
                    if lo < min_side or (lo and hi / lo > max_aspect):
                        continue          # a banner strip or sliver, not a product
                if dedupe_repeats and im.xref in repeated:
                    continue
                try:
                    d = doc.extract_image(im.xref)
                except Exception:
                    continue
                data = d["image"]
                product = _associate(im.bbox, items) if associate else None

                stem = _safe(product) if product else f"page{p.index + 1:03d}"
                count = used_names.get(stem, 0) + 1
                used_names[stem] = count
                fname = f"{stem}.{d['ext']}" if count == 1 else f"{stem}_{count}.{d['ext']}"
                with open(os.path.join(out_dir, fname), "wb") as f:
                    f.write(data)

                manifest.append({
                    "page": p.index + 1,
                    "product": product,
                    "collection": collection if product else None,
                    "file": fname,
                    "format": d["ext"],
                    "width_px": d["width"],
                    "height_px": d["height"],
                    "colorspace": _CS.get(d.get("colorspace"), d.get("colorspace")),
                    "bytes": len(data),
                    "dpi": _dpi(d["width"], im.bbox),
                    "dominant_color": _dominant_color(data),
                    "bbox": tuple(round(v, 1) for v in im.bbox) if im.bbox else None,
                })
    finally:
        doc.close()
    return manifest


def _pptx_title(slide) -> Optional[str]:
    try:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    return None


def export_pptx_images(
    path: str,
    out_dir: str,
    ocr: bool = False,
    only_photos: bool = True,
    min_side: int = 180,
    max_aspect: float = 4.0,
) -> List[dict]:
    """Extract product pictures from a PPTX deck into ``out_dir``; return a
    manifest. Each picture becomes a row anchor. Names come from the slide's
    text (title/text boxes); when a slide is pure image, ``ocr`` reads any text
    baked into the picture itself. Mirrors ``export_images`` for PDFs."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image

    from .layout import ocr as _ocr

    os.makedirs(out_dir, exist_ok=True)
    prs = Presentation(path)
    manifest: List[dict] = []
    used_names: dict = {}
    for i, slide in enumerate(prs.slides):
        title = _pptx_title(slide)
        slide_text = [sh.text_frame.text.strip() for sh in slide.shapes
                      if sh.has_text_frame and sh.text_frame.text.strip()]
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = shape.image.blob
                with Image.open(io.BytesIO(blob)) as im:
                    w, h = im.size
                    ext = (im.format or "png").lower()
            except Exception:
                continue
            if only_photos:
                lo, hi = min(w, h), max(w, h)
                if lo < min_side or (lo and hi / lo > max_aspect):
                    continue
            # Real text (title/text boxes) names the product; OCR of the picture is
            # kept separate because it only reliably yields codes/dimensions, not names.
            ocr_text = ""
            if ocr and not slide_text:
                ocr_text = " \n".join(wd.text for wd in _ocr.ocr_image_words(blob))
            stem = _safe(title) if title else f"slide{i + 1:03d}"
            count = used_names.get(stem, 0) + 1
            used_names[stem] = count
            fname = f"{stem}.{ext}" if count == 1 else f"{stem}_{count}.{ext}"
            with open(os.path.join(out_dir, fname), "wb") as f:
                f.write(blob)
            manifest.append({
                "page": i + 1,
                "title": title,
                "file": fname,
                "format": ext,
                "width_px": w,
                "height_px": h,
                "slide_text": " \n".join(slide_text),
                "ocr_text": ocr_text,
            })
    return manifest
