"""PPTX adapter (python-pptx).

PowerPoint is a gift compared to PDF: shapes carry positions, tables are native
(rows and cells, no reconstruction), and pictures expose real pixels. So we hand
tables straight through as high-confidence ``native_tables`` and only fall back
to geometry for free-floating text boxes.
"""
from __future__ import annotations

import io
from typing import List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..layout.images import classify
from ..model.elements import ImageRef, Page, Word
from ..model.table import ExtractedTable
from .base import DocumentAdapter, register

_EMU_PER_PT = 12700.0


def _pt(emu) -> float:
    return float(emu) / _EMU_PER_PT if emu is not None else 0.0


def _tokens_in_box(text: str, x0: float, y: float, x1: float, size: float, bold: bool) -> List[Word]:
    if not text.strip():
        return []
    width = max(1.0, x1 - x0)
    n = len(text)
    out: List[Word] = []
    idx = 0
    for tok in text.split():
        start = text.index(tok, idx)
        end = start + len(tok)
        idx = end
        wx0 = x0 + width * (start / n)
        wx1 = x0 + width * (end / n)
        out.append(Word(text=tok, x0=wx0, y0=y, x1=wx1, y1=y + max(size, 8), size=size, bold=bold))
    return out


def _slide_title(slide) -> Optional[str]:
    try:
        if slide.shapes.title is not None and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    return None


def _table_to_extracted(shape, page_index: int, title: Optional[str]) -> ExtractedTable:
    tbl = shape.table
    rows = list(tbl.rows)
    grid = [[cell.text.strip() for cell in row.cells] for row in rows]
    if not grid:
        return ExtractedTable(title, [], [], [], page_index, kind="grid")
    col_headers = grid[0]
    body = grid[1:] if len(grid) > 1 else []
    row_headers = [r[0] if r else "" for r in body]
    cells = [[(v if v else None) for v in (r[1:] if len(r) > 1 else r)] for r in body]
    col_headers = col_headers[1:] if len(col_headers) > 1 else col_headers
    return ExtractedTable(
        title=title,
        row_headers=row_headers,
        col_headers=col_headers,
        cells=cells if cells else [[v or None for v in r] for r in body],
        source_page=page_index,
        kind="grid",
    )


@register
class PptxAdapter(DocumentAdapter):
    extensions = ("pptx",)

    def parse(self, path: str) -> List[Page]:
        prs = Presentation(path)
        sw, sh = _pt(prs.slide_width), _pt(prs.slide_height)
        pages: List[Page] = []
        for i, slide in enumerate(prs.slides):
            words: List[Word] = []
            images: List[ImageRef] = []
            native: List[ExtractedTable] = []
            title = _slide_title(slide)
            for shape in slide.shapes:
                x0, y0 = _pt(shape.left), _pt(shape.top)
                x1 = x0 + _pt(shape.width)
                if shape.has_table:
                    native.append(_table_to_extracted(shape, i, title))
                    continue
                if shape.has_text_frame:
                    y = y0
                    for para in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in para.runs) or para.text
                        size = 0.0
                        bold = False
                        for run in para.runs:
                            if run.font.size is not None:
                                size = run.font.size.pt
                            if run.font.bold:
                                bold = True
                        words.extend(_tokens_in_box(text, x0, y, x1, size, bold))
                        y += max(size, 14) * 1.2
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        blob = shape.image.blob
                        from PIL import Image

                        with Image.open(io.BytesIO(blob)) as im:
                            w, h = im.size
                        images.append(ImageRef(page=i, width=w, height=h, kind=classify(w, h)))
                    except Exception:
                        pass
            pages.append(
                Page(
                    index=i,
                    width=sw,
                    height=sh,
                    words=words,
                    images=images,
                    source_kind="pptx",
                    native_tables=native,
                )
            )
        return pages
